from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

OUTPUT = Path(__file__).with_name("njangi_presentation_public.pptx")
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Njangi Digital"
slide.placeholders[1].text = "Solution de gestion de tontine numérique"

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Problème résolu"
body = slide.placeholders[1].text_frame
body.clear()
p = body.paragraphs[0]
p.text = "Les tontines traditionnelles manquent souvent de :"
for item in ["traçabilité", "transparence", "suivi des paiements", "gestion efficace des prêts"]:
    p = body.add_paragraph()
    p.text = f"• {item}"
    p.level = 0

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Ce que propose Njangi"
body = slide.placeholders[1].text_frame
body.clear()
for item in ["Gestion des groupes et membres", "Séances de cotisation et présences", "Prêts, remboursements et fond commun", "Rapports PDF et suivi financier"]:
    paragraph = body.add_paragraph()
    paragraph.text = f"• {item}"
    paragraph.level = 0

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Pourquoi c’est crédible"
body = slide.placeholders[1].text_frame
body.clear()
for item in ["Application déjà fonctionnelle", "Fonctionnalités essentielles complètes", "Tests validés", "Prête pour une présentation publique"]:
    paragraph = body.add_paragraph()
    paragraph.text = f"• {item}"
    paragraph.level = 0

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusion"
body = slide.placeholders[1].text_frame
body.clear()
body.paragraphs[0].text = "Njangi Digital est une solution numérique fiable, moderne et déjà opérationnelle pour la gestion de tontines."

prs.save(str(OUTPUT))
print(f"PowerPoint generated: {OUTPUT}")
